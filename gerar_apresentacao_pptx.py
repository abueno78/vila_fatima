"""
Gera apresentação PowerPoint bilíngue (PT/EN) da Unidade de Saúde Vila Fátima
Compatível com Microsoft PowerPoint e Google Apresentações (via importação)
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# Caminhos das fotos
FOTOS = {
    'aerea':      'fotos_unidade/foto_territorio_aereo.jpg',
    'odonto':     'fotos_unidade/foto_odontologia_pucrs.jpg',
    'acs_idosa':  'fotos_unidade/foto_acs_idosa.jpg',
    'equipe':     'fotos_unidade/foto_equipe_territorio.jpg',
    'graficos': {
        'pyramid':   'graficos_analise/pyramid_census.png',
        'raca':      'graficos_analise/raca_census.png',
        'vol':       'graficos_analise/4_pessoas_unicas_atendidas.png',
        'prof':      'graficos_analise/1_agrupamento_profissionais.png',
        'age_pct':   'graficos_analise/2_faixa_etaria_mensal_percentual.png',
        'ret_age':   'graficos_analise/7_retorno_faixas_gerais.png',
        'eld_donut': 'graficos_analise/5_distribuicao_idosos_geral.png',
        'eld_ret':   'graficos_analise/8_retorno_faixas_idosos.png',
    }
}

# Cores
PUCRS_BLUE  = RGBColor(0x00, 0x50, 0x8B)
PUCRS_GOLD  = RGBColor(0xEB, 0xB7, 0x00)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
DARK_BG     = RGBColor(0x0B, 0x0F, 0x19)
LIGHT_GRAY  = RGBColor(0x94, 0xA3, 0xB8)
DARK_CARD   = RGBColor(0x13, 0x1A, 0x2D)
GREEN_ACCENT= RGBColor(0x00, 0xDF, 0x89)

# Dimensões slide 16:9
W = Inches(13.333)
H = Inches(7.5)


def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_slide(prs):
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)


def bg(slide, color=DARK_BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, color, alpha=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text(slide, text, x, y, w, h, size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txb


def add_image_safe(slide, path, x, y, w, h):
    if os.path.exists(path):
        slide.shapes.add_picture(path, x, y, w, h)


def header_bar(slide, title_text):
    """Barra azul no topo com título."""
    add_rect(slide, 0, 0, W, Inches(0.85), PUCRS_BLUE)
    add_rect(slide, 0, Inches(0.85), W, Inches(0.06), PUCRS_GOLD)
    add_text(slide, title_text, Inches(0.35), Inches(0.15),
             Inches(11), Inches(0.7), size=22, bold=True, color=WHITE)


def footer(slide, lang='pt'):
    txt = {
        'pt': 'Unidade de Saúde Vila Fátima  •  APS PUCRS / SMS Porto Alegre',
        'en': 'Vila Fátima Primary Health Unit  •  APS PUCRS / SMS Porto Alegre'
    }[lang]
    add_rect(slide, 0, H - Inches(0.4), W, Inches(0.4), RGBColor(0x08, 0x0B, 0x14))
    add_text(slide, txt, Inches(0.35), H - Inches(0.38),
             Inches(12), Inches(0.35), size=9, color=LIGHT_GRAY)


def kpi_card(slide, label, value, x, y, w=Inches(3.8), h=Inches(1.3)):
    add_rect(slide, x, y, w, h, DARK_CARD)
    add_rect(slide, x, y, Inches(0.05), h, PUCRS_BLUE)
    add_text(slide, label, x + Inches(0.15), y + Inches(0.08),
             w - Inches(0.2), Inches(0.35), size=9, color=LIGHT_GRAY)
    add_text(slide, value, x + Inches(0.15), y + Inches(0.42),
             w - Inches(0.2), Inches(0.65), size=20, bold=True, color=RGBColor(0x00, 0x9E, 0xDB))


def bullet(slide, text, x, y, w, size=12, color=WHITE):
    add_text(slide, text, x, y, w, Inches(0.5), size=size, color=color, wrap=True)


def build(lang='pt'):
    prs = new_prs()

    t = {
        'pt': {
            'cover_title': 'Unidade de Saúde Vila Fátima',
            'cover_sub':   'Perfil Operacional e Demográfico do Território\nCenso 2022 & Atendimentos 2025–2026',
            'cover_tag':   'PUCRS Escola de Medicina  •  Porto Alegre, RS, Brasil',
            's1': '1. Contexto Territorial e Socioeconômico',
            's1_k': [('POPULAÇÃO RESIDENTE','4.891 pessoas'),
                     ('ÁREA TOTAL','0,5135 km²'),
                     ('COMUNIDADE URBANA','69,18% da área'),
                     ('RENDA MÉDIA MENSAL','R$ 1.721,14'),
                     ('RENDA MEDIANA MENSAL','R$ 1.212,00'),
                     ('DESVIO PADRÃO RENDA','R$ 378,07')],
            's1_b': ['• Alta densidade demográfica em território de alta vulnerabilidade social.',
                     '• 69,18% da área coberta por favelas e comunidades urbanas.',
                     '• Renda mediana equivalente a 1 salário mínimo (2022).'],
            's2': '2. Estrutura da População Residente (Pirâmide Etária)',
            's2_b': ['• Base larga: expressivo grupo de crianças e jovens (0-24 anos).',
                     '• Adultos consolidados: forte presença de 30 a 49 anos.',
                     '• Envelhecimento crescente: coorte de idosos (60+) em expansão.'],
            's3': '3. Composição por Raça e Cor (Censo 2022)',
            's3_b': ['• Brancos: ~43% | Pretos: ~32% | Pardos: ~24%.',
                     '• Mais de 55% da população é negra (pretos + pardos).',
                     '• Proporção superior à média municipal de Porto Alegre.'],
            's4': '4. Volume de Atendimentos e Alcance (2025–2026)',
            's4_b': ['• Fluxo mensal estável e contínuo de consultas.',
                     '• Alta taxa de retorno: pacientes acompanhados longitudinalmente.',
                     '• Vínculo forte entre equipe de saúde e famílias do território.'],
            's5': '5. Integração Multidisciplinar da Equipe',
            's5_b': ['• Medicina, Enfermagem, Odontologia e Técnicos integrados.',
                     '• Enfermagem protagoniza triagem, pré-natal e crônicos.',
                     '• Saúde bucal plenamente integrada ao plano de cuidados.'],
            's6': '6. Atendimentos por Faixa Etária (Evolução Mensal)',
            's6_b': ['• Adultos (30-59 anos) lideram o volume de consultas.',
                     '• Idosos (60+) mantêm alta frequência relativa todos os meses.',
                     '• Distribuição estável reflete agenda programática organizada.'],
            's7': '7. Frequência Média de Consultas por Faixa Etária',
            's7_b': ['• Idosos têm maior taxa de retorno — alta cronicidade.',
                     '• Crianças mantêm visitas regulares (puericultura, vacinas).',
                     '• Alta taxa confirma acompanhamento longitudinal efetivo.'],
            's8': '8. Perfil Geriátrico — Coorte 60+',
            's8_b': ['• Grupo 60-70 anos: maior fatia absoluta de idosos.',
                     '• Grupo 81-90 anos: média de 17,42 visitas por paciente.',
                     '• Alta complexidade clínica exige retornos frequentes.'],
            's9': '9. Conclusões e Impacto Institucional',
            's9_b': ['• Modelo territorializado de saúde da família com alta efetividade.',
                     '• Cooperação PUCRS–SMS eleva qualidade da APS local.',
                     '• Gestão baseada em dados otimiza alocação de recursos.',
                     '• Alta taxa de retorno confirma vínculo longitudinal robusto.',
                     '• Vila Fátima: referência pedagógica e clínica de excelência.'],
        },
        'en': {
            'cover_title': 'Vila Fátima Primary Health Unit',
            'cover_sub':   'Territorial Operational & Demographic Profile\nCensus 2022 & 2025–2026 Consultation Data',
            'cover_tag':   'PUCRS School of Medicine  •  Porto Alegre, RS, Brazil',
            's1': '1. Territorial & Socioeconomic Context',
            's1_k': [('RESIDENT POPULATION','4,891 residents'),
                     ('TOTAL AREA','0.5135 km²'),
                     ('URBAN COMMUNITY','69.18% of area'),
                     ('MEAN MONTHLY INCOME','R$ 1,721.14'),
                     ('MEDIAN MONTHLY INCOME','R$ 1,212.00'),
                     ('INCOME STD. DEVIATION','R$ 378.07')],
            's1_b': ['• High demographic density in a high social vulnerability area.',
                     '• 69.18% of the territory covered by informal communities (favelas).',
                     '• Median income equivalent to 1 minimum wage (2022).'],
            's2': '2. Resident Population Structure (Population Pyramid)',
            's2_b': ['• Broad base: significant children and youth cohort (0-24 years).',
                     '• Consolidated adults: strong presence aged 30-49.',
                     '• Growing elderly cohort (60+) requiring chronic disease management.'],
            's3': '3. Racial & Ethnic Composition (Census 2022)',
            's3_b': ['• White: ~43% | Black: ~32% | Mixed (Pardo): ~24%.',
                     '• Over 55% of residents self-declare as Black or Mixed.',
                     '• Proportion higher than Porto Alegre municipal average.'],
            's4': '4. Consultation Volume & Reach (2025–2026)',
            's4_b': ['• Stable and continuous monthly consultation flow.',
                     '• High return rate: patients monitored longitudinally.',
                     '• Strong bonding between health team and community families.'],
            's5': '5. Multidisciplinary Team Integration',
            's5_b': ['• Medicine, Nursing, Dentistry and Technicians integrated.',
                     '• Nursing leads triage, prenatal care and chronic disease control.',
                     '• Oral health fully integrated into the care plan.'],
            's6': '6. Consultations by Age Group (Monthly Trend)',
            's6_b': ['• Adults (30-59) lead monthly consultation volumes.',
                     '• Elderly (60+) maintain high relative frequency every month.',
                     '• Stable distribution reflects organized programmatic scheduling.'],
            's7': '7. Average Visit Frequency by Age Group',
            's7_b': ['• Elderly have highest return rate — high chronicity.',
                     '• Children maintain regular visits (pediatrics, vaccines).',
                     '• High rate confirms effective longitudinal follow-up.'],
            's8': '8. Geriatric Care Profile — 60+ Cohort',
            's8_b': ['• 60-70 cohort: largest absolute share of elderly patients.',
                     '• 81-90 cohort: average of 17.42 visits per patient.',
                     '• High clinical complexity demands frequent returns.'],
            's9': '9. Strategic & Institutional Conclusions',
            's9_b': ['• Highly effective territorial family health care model.',
                     '• PUCRS–SMS cooperation elevates local primary care quality.',
                     '• Data-driven management optimizes resource allocation.',
                     '• High return rate confirms robust longitudinal patient bonding.',
                     '• Vila Fátima: a reference for academic and clinical excellence.'],
        }
    }[lang]

    # =========================================================================
    # SLIDE 1 — CAPA
    # =========================================================================
    sl = blank_slide(prs)
    bg(sl, DARK_BG)

    # Foto aérea como background com overlay
    if os.path.exists(FOTOS['aerea']):
        sl.shapes.add_picture(FOTOS['aerea'], 0, 0, W, H)
    # Overlay escuro
    add_rect(sl, 0, 0, W, H, RGBColor(0x05, 0x08, 0x12))
    # Faixa azul no lado esquerdo
    add_rect(sl, 0, 0, Inches(0.5), H, PUCRS_BLUE)
    add_rect(sl, Inches(0.5), 0, Inches(0.07), H, PUCRS_GOLD)

    # Logo PUCRS (se disponível localmente)
    if os.path.exists('logo_pucrs.png'):
        sl.shapes.add_picture('logo_pucrs.png', Inches(1.0), Inches(1.2), Inches(2.5), Inches(0.9))

    add_text(sl, t['cover_title'], Inches(1.0), Inches(2.4), Inches(10),
             Inches(1.2), size=38, bold=True, color=WHITE)
    add_text(sl, t['cover_sub'], Inches(1.0), Inches(3.7), Inches(10),
             Inches(1.5), size=18, color=LIGHT_GRAY, wrap=True)
    add_rect(sl, Inches(1.0), Inches(5.1), Inches(3), Inches(0.06), PUCRS_GOLD)
    add_text(sl, t['cover_tag'], Inches(1.0), Inches(5.3), Inches(10),
             Inches(0.5), size=12, color=LIGHT_GRAY)

    footer(sl, lang)

    # =========================================================================
    # SLIDE 2 — CONTEXTO TERRITORIAL
    # =========================================================================
    sl = blank_slide(prs)
    bg(sl)
    header_bar(sl, t['s1'])

    # 6 KPI cards (2 linhas x 3 colunas)
    cols = [Inches(0.3), Inches(4.5), Inches(8.7)]
    rows = [Inches(1.1), Inches(2.65)]
    for i, (lbl, val) in enumerate(t['s1_k']):
        col = cols[i % 3]
        row = rows[i // 3]
        kpi_card(sl, lbl, val, col, row, Inches(3.9), Inches(1.35))

    # Bullets abaixo
    for i, b_txt in enumerate(t['s1_b']):
        add_text(sl, b_txt, Inches(0.3), Inches(4.25) + i * Inches(0.42),
                 Inches(7), Inches(0.4), size=11, color=LIGHT_GRAY, wrap=True)

    # Foto aérea à direita
    add_image_safe(sl, FOTOS['aerea'], Inches(7.5), Inches(4.0), Inches(5.5), Inches(3.1))

    footer(sl, lang)

    # =========================================================================
    # SLIDE 3 — PIRÂMIDE ETÁRIA
    # =========================================================================
    sl = blank_slide(prs)
    bg(sl)
    header_bar(sl, t['s2'])
    add_image_safe(sl, FOTOS['graficos']['pyramid'], Inches(0.2), Inches(1.0), Inches(7.5), Inches(5.5))
    add_rect(sl, Inches(7.9), Inches(1.0), Inches(5.0), Inches(5.5), DARK_CARD)
    for i, b_txt in enumerate(t['s2_b']):
        add_text(sl, b_txt, Inches(8.1), Inches(1.3) + i * Inches(0.8),
                 Inches(4.7), Inches(0.75), size=12, color=WHITE, wrap=True)
    footer(sl, lang)

    # =========================================================================
    # SLIDE 4 — RAÇA/COR
    # =========================================================================
    sl = blank_slide(prs)
    bg(sl)
    header_bar(sl, t['s3'])
    add_image_safe(sl, FOTOS['graficos']['raca'], Inches(0.2), Inches(1.0), Inches(7.5), Inches(5.5))
    add_rect(sl, Inches(7.9), Inches(1.0), Inches(5.0), Inches(5.5), DARK_CARD)
    for i, b_txt in enumerate(t['s3_b']):
        add_text(sl, b_txt, Inches(8.1), Inches(1.3) + i * Inches(0.8),
                 Inches(4.7), Inches(0.75), size=12, color=WHITE, wrap=True)
    footer(sl, lang)

    # =========================================================================
    # SLIDE 5 — VOLUME DE ATENDIMENTOS
    # =========================================================================
    sl = blank_slide(prs)
    bg(sl)
    header_bar(sl, t['s4'])
    add_image_safe(sl, FOTOS['graficos']['vol'], Inches(0.2), Inches(1.0), Inches(7.5), Inches(4.0))
    add_rect(sl, Inches(7.9), Inches(1.0), Inches(5.0), Inches(4.0), DARK_CARD)
    for i, b_txt in enumerate(t['s4_b']):
        add_text(sl, b_txt, Inches(8.1), Inches(1.3) + i * Inches(0.8),
                 Inches(4.7), Inches(0.75), size=12, color=WHITE, wrap=True)
    # Foto equipe no território embaixo
    add_image_safe(sl, FOTOS['equipe'], Inches(0.2), Inches(5.1), Inches(12.9), Inches(2.0))
    footer(sl, lang)

    # =========================================================================
    # SLIDE 6 — PROFISSIONAIS (com foto odontologia)
    # =========================================================================
    sl = blank_slide(prs)
    bg(sl)
    header_bar(sl, t['s5'])
    add_image_safe(sl, FOTOS['graficos']['prof'], Inches(0.2), Inches(1.0), Inches(6.5), Inches(4.5))
    # Foto da dentista PUCRS
    add_image_safe(sl, FOTOS['odonto'], Inches(6.8), Inches(1.0), Inches(3.0), Inches(4.5))
    add_rect(sl, Inches(10.0), Inches(1.0), Inches(3.1), Inches(4.5), DARK_CARD)
    for i, b_txt in enumerate(t['s5_b']):
        add_text(sl, b_txt, Inches(10.15), Inches(1.3) + i * Inches(0.9),
                 Inches(2.9), Inches(0.85), size=11, color=WHITE, wrap=True)
    footer(sl, lang)

    # =========================================================================
    # SLIDE 7 — EVOLUÇÃO POR FAIXA ETÁRIA
    # =========================================================================
    sl = blank_slide(prs)
    bg(sl)
    header_bar(sl, t['s6'])
    add_image_safe(sl, FOTOS['graficos']['age_pct'], Inches(0.2), Inches(1.0), Inches(8.5), Inches(5.5))
    add_rect(sl, Inches(8.9), Inches(1.0), Inches(4.1), Inches(5.5), DARK_CARD)
    for i, b_txt in enumerate(t['s6_b']):
        add_text(sl, b_txt, Inches(9.1), Inches(1.3) + i * Inches(0.8),
                 Inches(3.8), Inches(0.75), size=11, color=WHITE, wrap=True)
    footer(sl, lang)

    # =========================================================================
    # SLIDE 8 — TAXA DE RETORNO
    # =========================================================================
    sl = blank_slide(prs)
    bg(sl)
    header_bar(sl, t['s7'])
    add_image_safe(sl, FOTOS['graficos']['ret_age'], Inches(0.2), Inches(1.0), Inches(8.5), Inches(5.5))
    add_rect(sl, Inches(8.9), Inches(1.0), Inches(4.1), Inches(5.5), DARK_CARD)
    for i, b_txt in enumerate(t['s7_b']):
        add_text(sl, b_txt, Inches(9.1), Inches(1.3) + i * Inches(0.8),
                 Inches(3.8), Inches(0.75), size=11, color=WHITE, wrap=True)
    footer(sl, lang)

    # =========================================================================
    # SLIDE 9 — GERIATRIA (com foto ACS e idosa)
    # =========================================================================
    sl = blank_slide(prs)
    bg(sl)
    header_bar(sl, t['s8'])
    add_image_safe(sl, FOTOS['graficos']['eld_donut'], Inches(0.2), Inches(1.0), Inches(4.0), Inches(4.0))
    add_image_safe(sl, FOTOS['graficos']['eld_ret'], Inches(4.4), Inches(1.0), Inches(5.0), Inches(4.0))
    # Foto ACS com idosa
    add_image_safe(sl, FOTOS['acs_idosa'], Inches(9.6), Inches(1.0), Inches(3.5), Inches(4.0))
    add_rect(sl, Inches(0.2), Inches(5.15), Inches(12.9), Inches(1.85), DARK_CARD)
    for i, b_txt in enumerate(t['s8_b']):
        add_text(sl, b_txt, Inches(0.4), Inches(5.3) + i * Inches(0.55),
                 Inches(12.5), Inches(0.5), size=11, color=WHITE, wrap=True)
    footer(sl, lang)

    # =========================================================================
    # SLIDE 10 — CONCLUSÕES
    # =========================================================================
    sl = blank_slide(prs)
    bg(sl)
    header_bar(sl, t['s9'])

    # Duas colunas de cards
    left_x  = Inches(0.3)
    right_x = Inches(6.8)
    add_rect(sl, left_x,  Inches(1.05), Inches(6.3), Inches(5.7), DARK_CARD)
    add_rect(sl, right_x, Inches(1.05), Inches(6.2), Inches(5.7), DARK_CARD)
    add_rect(sl, left_x,  Inches(1.05), Inches(0.07), Inches(5.7), PUCRS_BLUE)
    add_rect(sl, right_x, Inches(1.05), Inches(0.07), Inches(5.7), GREEN_ACCENT)

    # Esquerda: bullets 0-2
    for i, b_txt in enumerate(t['s9_b'][:3]):
        add_text(sl, b_txt, left_x + Inches(0.2), Inches(1.4) + i * Inches(1.0),
                 Inches(5.9), Inches(0.9), size=12, color=WHITE, wrap=True)

    # Direita: bullets 3-4 + foto equipe
    for i, b_txt in enumerate(t['s9_b'][3:]):
        add_text(sl, b_txt, right_x + Inches(0.2), Inches(1.4) + i * Inches(1.0),
                 Inches(5.9), Inches(0.9), size=12, color=WHITE, wrap=True)
    add_image_safe(sl, FOTOS['equipe'], right_x + Inches(0.2), Inches(3.5), Inches(5.8), Inches(2.9))

    footer(sl, lang)

    out = f'apresentacao_vila_fatima_{lang}.pptx'
    prs.save(out)
    print(f'[OK] {out}')


if __name__ == '__main__':
    print('Gerando PowerPoint PT...')
    build('pt')
    print('Gerando PowerPoint EN...')
    build('en')
    print('Concluído!')
